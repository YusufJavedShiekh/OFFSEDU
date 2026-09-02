"""
StudyGemma - PPTX Processor

Responsibilities:
- Validate PPTX files
- Extract presentation metadata
- Process slides individually
- Extract titles and text
- Detect headings and bullet/list items
- Extract tables
- Detect images
- Detect shapes
- Extract speaker notes when available
- Optionally send images to OCR
- Clean extracted content
- Preserve slide order and structure
- Validate extraction quality
- Prepare RAG-ready data

Does NOT:
- Store files
- Generate embeddings
- Perform vector search
- Call Gemma
"""

from __future__ import annotations

import hashlib
import logging
import re
import zipfile
from pathlib import Path
from typing import Any, Optional, Union

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


logger = logging.getLogger(__name__)


class PPTXProcessor:
    """
    Processor for Microsoft PowerPoint PPTX presentations.
    """

    DEFAULT_MIN_TEXT_LENGTH = 20

    def __init__(
        self,
        min_text_length: int = DEFAULT_MIN_TEXT_LENGTH,
        ocr_service: Optional[Any] = None,
    ):
        self.min_text_length = max(
            1,
            min_text_length,
        )

        self.ocr_service = ocr_service

        if Presentation is None:
            logger.warning(
                "python-pptx is not installed. "
                "Install it with: pip install python-pptx"
            )

    # =========================================================
    # MAIN PROCESS
    # =========================================================

    def process(
        self,
        pptx_path: Union[str, Path],
        use_ocr: bool = True,
        include_notes: bool = True,
    ) -> dict[str, Any]:
        """
        Process a complete PPTX presentation.
        """

        try:

            path = Path(
                pptx_path
            ).resolve()

            validation = self.validate_file(
                path
            )

            if not validation["valid"]:
                return self._failure(
                    validation["error"]
                )

            if Presentation is None:
                return self._failure(
                    "python-pptx is not installed."
                )

            presentation = Presentation(
                str(path)
            )

            metadata = (
                self.extract_metadata(
                    presentation,
                    path,
                )
            )

            slides = []
            rag_slides = []

            for index, slide in enumerate(
                presentation.slides
            ):

                slide_result = (
                    self.process_slide(
                        slide=slide,
                        slide_number=index + 1,
                        use_ocr=use_ocr,
                        include_notes=include_notes,
                    )
                )

                slides.append(
                    slide_result
                )

                if slide_result.get(
                    "success",
                    False,
                ):

                    rag_slides.append(
                        self.create_rag_slide(
                            slide_result
                        )
                    )

            full_text = (
                self.combine_slides(
                    slides
                )
            )

            quality = (
                self.validate_quality(
                    slides=slides,
                    full_text=full_text,
                )
            )

            structure = (
                self.detect_presentation_structure(
                    slides
                )
            )

            return {
                "success": True,
                "document": metadata,
                "slide_count": len(
                    presentation.slides
                ),
                "slides": slides,
                "full_text": full_text,
                "structure": structure,
                "quality": quality,
                "rag_ready": {
                    "document_id": metadata[
                        "document_id"
                    ],
                    "source_type": "pptx",
                    "slide_count": len(
                        rag_slides
                    ),
                    "slides": rag_slides,
                    "text": full_text,
                },
            }

        except Exception as error:

            logger.exception(
                "PPTX processing failed."
            )

            return self._failure(
                str(error)
            )

    # =========================================================
    # VALIDATION
    # =========================================================

    @staticmethod
    def validate_file(
        path: Path,
    ) -> dict[str, Any]:
        """
        Validate the PPTX file and its ZIP structure.
        """

        if not path.exists():

            return {
                "valid": False,
                "error": (
                    f"PPTX file not found: {path}"
                ),
            }

        if not path.is_file():

            return {
                "valid": False,
                "error": (
                    "Provided path is not a file."
                ),
            }

        if path.suffix.lower() != ".pptx":

            return {
                "valid": False,
                "error": (
                    "Provided file is not a PPTX."
                ),
            }

        try:

            if not zipfile.is_zipfile(
                path
            ):

                return {
                    "valid": False,
                    "error": (
                        "File is not a valid PPTX "
                        "container."
                    ),
                }

            with zipfile.ZipFile(
                path,
                "r",
            ) as archive:

                names = archive.namelist()

                if (
                    "[Content_Types].xml"
                    not in names
                ):

                    return {
                        "valid": False,
                        "error": (
                            "Invalid PPTX structure."
                        ),
                    }

                if not any(
                    name.startswith(
                        "ppt/slides/slide"
                    )
                    for name in names
                ):

                    return {
                        "valid": False,
                        "error": (
                            "No presentation slides "
                            "were found."
                        ),
                    }

        except zipfile.BadZipFile:

            return {
                "valid": False,
                "error": (
                    "PPTX file is corrupted."
                ),
            }

        return {
            "valid": True,
            "error": None,
        }

    # =========================================================
    # METADATA
    # =========================================================

    @staticmethod
    def extract_metadata(
        presentation: Any,
        path: Path,
    ) -> dict[str, Any]:
        """
        Extract PowerPoint core metadata.
        """

        properties = (
            presentation.core_properties
        )

        return {
            "document_id": (
                PPTXProcessor.generate_document_id(
                    path
                )
            ),
            "filename": path.name,
            "path": str(path),
            "title": (
                properties.title
                or ""
            ),
            "subject": (
                properties.subject
                or ""
            ),
            "author": (
                properties.author
                or ""
            ),
            "keywords": (
                properties.keywords
                or ""
            ),
            "comments": (
                properties.comments
                or ""
            ),
            "category": (
                properties.category
                or ""
            ),
            "last_modified_by": (
                properties.last_modified_by
                or ""
            ),
            "created": (
                str(properties.created)
                if properties.created
                else None
            ),
            "modified": (
                str(properties.modified)
                if properties.modified
                else None
            ),
            "slide_count": len(
                presentation.slides
            ),
        }

    # =========================================================
    # PROCESS SLIDE
    # =========================================================

    def process_slide(
        self,
        slide: Any,
        slide_number: int,
        use_ocr: bool = True,
        include_notes: bool = True,
    ) -> dict[str, Any]:
        """
        Process one PowerPoint slide.
        """

        try:

            title = self.extract_slide_title(
                slide
            )

            text_blocks = (
                self.extract_text_blocks(
                    slide
                )
            )

            tables = (
                self.extract_tables(
                    slide
                )
            )

            images = (
                self.extract_images(
                    slide
                )
            )

            shapes = (
                self.extract_shapes(
                    slide
                )
            )

            notes = ""

            if include_notes:

                notes = (
                    self.extract_notes(
                        slide
                    )
                )

            ocr_results = []

            if (
                use_ocr
                and self.ocr_service is not None
                and images
            ):

                ocr_results = (
                    self.process_images_with_ocr(
                        slide
                    )
                )

            text = (
                self.build_slide_text(
                    title=title,
                    text_blocks=text_blocks,
                    tables=tables,
                    notes=notes,
                    ocr_results=ocr_results,
                )
            )

            cleaned_text = (
                self.clean_text(
                    text
                )
            )

            slide_type = (
                self.detect_slide_type(
                    title=title,
                    text=cleaned_text,
                    image_count=len(images),
                    table_count=len(tables),
                    shape_count=len(shapes),
                    ocr_used=bool(
                        ocr_results
                    ),
                )
            )

            structure = (
                self.detect_slide_structure(
                    title=title,
                    text_blocks=text_blocks,
                    tables=tables,
                    images=images,
                    shapes=shapes,
                    notes=notes,
                )
            )

            return {
                "success": True,
                "slide_number": slide_number,
                "title": title,
                "text": cleaned_text,
                "text_length": len(
                    cleaned_text
                ),
                "text_blocks": text_blocks,
                "tables": tables,
                "images": images,
                "shapes": shapes,
                "notes": notes,
                "ocr_results": ocr_results,
                "ocr_used": bool(
                    ocr_results
                ),
                "slide_type": slide_type,
                "structure": structure,
            }

        except Exception as error:

            logger.exception(
                "Failed to process slide %s.",
                slide_number,
            )

            return {
                "success": False,
                "slide_number": slide_number,
                "error": str(error),
            }

    # =========================================================
    # SLIDE TITLE
    # =========================================================

    @staticmethod
    def extract_slide_title(
        slide: Any,
    ) -> str:
        """
        Extract the slide title using PowerPoint's
        title placeholder when available.
        """

        try:

            if slide.shapes.title:

                title = (
                    slide.shapes.title.text
                )

                if title:

                    return PPTXProcessor.clean_text(
                        title
                    )

        except Exception:
            pass

        # -----------------------------------------------------
        # Fallback:
        # Find a short text block near the top.
        # -----------------------------------------------------

        candidates = []

        for shape in slide.shapes:

            if not getattr(
                shape,
                "has_text_frame",
                False,
            ):
                continue

            text = (
                shape.text
                or ""
            ).strip()

            if not text:
                continue

            try:

                top = shape.top

            except Exception:

                top = 0

            candidates.append(
                (
                    top,
                    text,
                )
            )

        candidates.sort(
            key=lambda item: item[0]
        )

        for _, text in candidates:

            if len(text) <= 150:

                return PPTXProcessor.clean_text(
                    text
                )

        return ""

    # =========================================================
    # TEXT BLOCKS
    # =========================================================

    @staticmethod
    def extract_text_blocks(
        slide: Any,
    ) -> list[dict[str, Any]]:
        """
        Extract text from slide shapes while
        preserving shape order.
        """

        blocks = []

        for index, shape in enumerate(
            slide.shapes
        ):

            if not getattr(
                shape,
                "has_text_frame",
                False,
            ):
                continue

            text = (
                PPTXProcessor.clean_text(
                    shape.text
                )
            )

            if not text:
                continue

            is_title = False

            try:

                is_title = (
                    shape == slide.shapes.title
                )

            except Exception:
                pass

            paragraphs = []

            try:

                for paragraph in (
                    shape.text_frame.paragraphs
                ):

                    paragraph_text = (
                        PPTXProcessor.clean_text(
                            paragraph.text
                        )
                    )

                    if not paragraph_text:
                        continue

                    level = getattr(
                        paragraph,
                        "level",
                        0,
                    )

                    paragraphs.append(
                        {
                            "text": paragraph_text,
                            "level": level,
                            "is_bullet": (
                                PPTXProcessor.is_bullet(
                                    paragraph
                                )
                            ),
                        }
                    )

            except Exception:
                pass

            blocks.append(
                {
                    "shape_index": index,
                    "text": text,
                    "is_title": is_title,
                    "paragraphs": paragraphs,
                }
            )

        return blocks

    # =========================================================
    # BULLET DETECTION
    # =========================================================

    @staticmethod
    def is_bullet(
        paragraph: Any,
    ) -> bool:
        """
        Detect whether a PowerPoint paragraph
        is formatted as a bullet.
        """

        text = (
            getattr(
                paragraph,
                "text",
                "",
            )
            or ""
        ).strip()

        if re.match(
            r"^[-*•]\s+",
            text,
        ):

            return True

        # PowerPoint bullet formatting is represented
        # internally and is not always exposed uniformly
        # by python-pptx. Level > 0 is therefore a useful
        # supporting signal, but not definitive.
        level = getattr(
            paragraph,
            "level",
            0,
        )

        return level > 0

    # =========================================================
    # TABLE EXTRACTION
    # =========================================================

    @staticmethod
    def extract_tables(
        slide: Any,
    ) -> list[dict[str, Any]]:
        """
        Extract all PowerPoint tables.
        """

        tables = []

        for shape_index, shape in enumerate(
            slide.shapes
        ):

            if not getattr(
                shape,
                "has_table",
                False,
            ):

                continue

            table = shape.table

            rows = []

            for row_index, row in enumerate(
                table.rows
            ):

                cells = []

                for column_index, cell in enumerate(
                    row.cells
                ):

                    cells.append(
                        {
                            "column": column_index,
                            "text": (
                                PPTXProcessor.clean_text(
                                    cell.text
                                )
                            ),
                        }
                    )

                rows.append(
                    {
                        "row": row_index,
                        "cells": cells,
                    }
                )

            tables.append(
                {
                    "shape_index": shape_index,
                    "rows": rows,
                    "row_count": len(
                        rows
                    ),
                    "column_count": (
                        len(
                            table.columns
                        )
                    ),
                }
            )

        return tables

    # =========================================================
    # IMAGE EXTRACTION
    # =========================================================

    @staticmethod
    def extract_images(
        slide: Any,
    ) -> list[dict[str, Any]]:
        """
        Extract image metadata.

        Image bytes are not saved here.
        """

        images = []

        for index, shape in enumerate(
            slide.shapes
        ):

            if (
                not hasattr(
                    shape,
                    "image",
                )
            ):

                continue

            try:

                image = shape.image

                images.append(
                    {
                        "shape_index": index,
                        "filename": (
                            image.filename
                        ),
                        "extension": (
                            image.ext
                        ),
                        "content_type": (
                            image.content_type
                        ),
                        "width": (
                            shape.width
                        ),
                        "height": (
                            shape.height
                        ),
                    }
                )

            except Exception as error:

                logger.warning(
                    "Could not read PPTX image: %s",
                    error,
                )

        return images

    # =========================================================
    # SHAPE EXTRACTION
    # =========================================================

    @staticmethod
    def extract_shapes(
        slide: Any,
    ) -> list[dict[str, Any]]:
        """
        Extract basic information about non-text,
        non-table shapes.

        This helps preserve diagrams and visual layout
        information for later processing.
        """

        shapes = []

        for index, shape in enumerate(
            slide.shapes
        ):

            shape_type = str(
                getattr(
                    shape,
                    "shape_type",
                    "",
                )
            )

            shape_info = {
                "shape_index": index,
                "shape_type": shape_type,
                "name": (
                    getattr(
                        shape,
                        "name",
                        "",
                    )
                    or ""
                ),
                "has_text": bool(
                    getattr(
                        shape,
                        "has_text_frame",
                        False,
                    )
                ),
                "has_table": bool(
                    getattr(
                        shape,
                        "has_table",
                        False,
                    )
                ),
                "left": getattr(
                    shape,
                    "left",
                    None,
                ),
                "top": getattr(
                    shape,
                    "top",
                    None,
                ),
                "width": getattr(
                    shape,
                    "width",
                    None,
                ),
                "height": getattr(
                    shape,
                    "height",
                    None,
                ),
            }

            shapes.append(
                shape_info
            )

        return shapes

    # =========================================================
    # SPEAKER NOTES
    # =========================================================

    @staticmethod
    def extract_notes(
        slide: Any,
    ) -> str:
        """
        Extract speaker notes when supported by
        the installed python-pptx version.
        """

        try:

            notes_slide = (
                slide.notes_slide
            )

            text_parts = []

            for shape in (
                notes_slide.shapes
            ):

                if not getattr(
                    shape,
                    "has_text_frame",
                    False,
                ):

                    continue

                text = (
                    shape.text
                    or ""
                ).strip()

                if text:
                    text_parts.append(
                        text
                    )

            return PPTXProcessor.clean_text(
                "\n".join(
                    text_parts
                )
            )

        except Exception as error:

            logger.debug(
                "Speaker notes unavailable: %s",
                error,
            )

            return ""

    # =========================================================
    # OCR IMAGES
    # =========================================================

    def process_images_with_ocr(
        self,
        slide: Any,
    ) -> list[dict[str, Any]]:
        """
        Send embedded slide images to OCR.

        Image bytes are processed in memory and are
        not permanently saved by this method.
        """

        if self.ocr_service is None:

            return []

        results = []

        for index, shape in enumerate(
            slide.shapes
        ):

            if not hasattr(
                shape,
                "image",
            ):

                continue

            try:

                image_bytes = (
                    shape.image.blob
                )

                text = self.run_ocr(
                    image_bytes
                )

                if text:

                    results.append(
                        {
                            "image_index": index + 1,
                            "filename": (
                                shape.image.filename
                            ),
                            "text": (
                                self.clean_text(
                                    text
                                )
                            ),
                        }
                    )

            except Exception as error:

                logger.warning(
                    "OCR failed for slide image %s: %s",
                    index,
                    error,
                )

        return results

    # =========================================================
    # OCR SERVICE
    # =========================================================

    def run_ocr(
        self,
        image_bytes: bytes,
    ) -> str:
        """
        Call the configured OCR service.
        """

        if self.ocr_service is None:

            return ""

        try:

            if hasattr(
                self.ocr_service,
                "extract_text",
            ):

                result = (
                    self.ocr_service.extract_text(
                        image_bytes
                    )
                )

            elif hasattr(
                self.ocr_service,
                "process_image",
            ):

                result = (
                    self.ocr_service.process_image(
                        image_bytes
                    )
                )

            elif callable(
                self.ocr_service
            ):

                result = (
                    self.ocr_service(
                        image_bytes
                    )
                )

            else:

                return ""

            if isinstance(
                result,
                str,
            ):

                return result

            if isinstance(
                result,
                dict,
            ):

                return str(
                    result.get(
                        "text",
                        "",
                    )
                )

            return str(
                result or ""
            )

        except Exception as error:

            logger.warning(
                "OCR service failed: %s",
                error,
            )

            return ""

    # =========================================================
    # BUILD SLIDE TEXT
    # =========================================================

    @staticmethod
    def build_slide_text(
        title: str,
        text_blocks: list[dict[str, Any]],
        tables: list[dict[str, Any]],
        notes: str,
        ocr_results: list[dict[str, Any]],
    ) -> str:
        """
        Build a unified textual representation of a slide.
        """

        sections = []

        # -----------------------------------------------------
        # Title
        # -----------------------------------------------------

        if title:

            sections.append(
                f"[TITLE]\n{title}"
            )

        # -----------------------------------------------------
        # Text
        # -----------------------------------------------------

        for block in text_blocks:

            text = block.get(
                "text",
                "",
            )

            if not text:
                continue

            if block.get(
                "is_title",
                False,
            ):

                continue

            sections.append(
                text
            )

        # -----------------------------------------------------
        # Tables
        # -----------------------------------------------------

        for table in tables:

            sections.append(
                "\n[TABLE]"
            )

            for row in table.get(
                "rows",
                [],
            ):

                values = [
                    cell.get(
                        "text",
                        "",
                    )
                    for cell in row.get(
                        "cells",
                        [],
                    )
                ]

                sections.append(
                    " | ".join(
                        values
                    )
                )

            sections.append(
                "[/TABLE]"
            )

        # -----------------------------------------------------
        # Speaker notes
        # -----------------------------------------------------

        if notes:

            sections.append(
                f"\n[NOTES]\n{notes}"
            )

        # -----------------------------------------------------
        # OCR
        # -----------------------------------------------------

        for result in ocr_results:

            text = result.get(
                "text",
                "",
            )

            if text:

                sections.append(
                    f"\n[IMAGE TEXT]\n{text}"
                )

        return "\n\n".join(
            sections
        )

    # =========================================================
    # SLIDE TYPE
    # =========================================================

    @staticmethod
    def detect_slide_type(
        title: str,
        text: str,
        image_count: int,
        table_count: int,
        shape_count: int,
        ocr_used: bool,
    ) -> str:
        """
        Basic structural slide classification.
        """

        if ocr_used:

            return "image_based"

        if (
            table_count > 0
            and not text
        ):

            return "table"

        if (
            image_count > 0
            and not text
        ):

            return "visual"

        if (
            image_count > 0
            and text
        ):

            return "mixed"

        if text:

            if title:

                return "text"

            return "text_without_title"

        if shape_count > 0:

            return "diagram_or_graphic"

        return "empty"

    # =========================================================
    # SLIDE STRUCTURE
    # =========================================================

    @staticmethod
    def detect_slide_structure(
        title: str,
        text_blocks: list[dict[str, Any]],
        tables: list[dict[str, Any]],
        images: list[dict[str, Any]],
        shapes: list[dict[str, Any]],
        notes: str,
    ) -> dict[str, Any]:
        """
        Produce structural information for one slide.
        """

        bullet_count = 0

        for block in text_blocks:

            for paragraph in block.get(
                "paragraphs",
                [],
            ):

                if paragraph.get(
                    "is_bullet",
                    False,
                ):

                    bullet_count += 1

        return {
            "has_title": bool(
                title
            ),
            "has_text": bool(
                text_blocks
            ),
            "has_bullets": (
                bullet_count > 0
            ),
            "bullet_count": bullet_count,
            "has_tables": bool(
                tables
            ),
            "has_images": bool(
                images
            ),
            "has_shapes": bool(
                shapes
            ),
            "has_notes": bool(
                notes
            ),
        }

    # =========================================================
    # PRESENTATION STRUCTURE
    # =========================================================

    @staticmethod
    def detect_presentation_structure(
        slides: list[dict[str, Any]],
    ) -> dict[str, Any]:
        """
        Summarize the structure of the complete presentation.
        """

        return {
            "slide_count": len(
                slides
            ),
            "slides_with_titles": sum(
                1
                for slide in slides
                if slide.get(
                    "title"
                )
            ),
            "slides_with_images": sum(
                1
                for slide in slides
                if slide.get(
                    "images"
                )
            ),
            "slides_with_tables": sum(
                1
                for slide in slides
                if slide.get(
                    "tables"
                )
            ),
            "slides_with_notes": sum(
                1
                for slide in slides
                if slide.get(
                    "notes"
                )
            ),
            "slides_with_ocr": sum(
                1
                for slide in slides
                if slide.get(
                    "ocr_used",
                    False,
                )
            ),
        }

    # =========================================================
    # COMBINE SLIDES
    # =========================================================

    @staticmethod
    def combine_slides(
        slides: list[dict[str, Any]],
    ) -> str:
        """
        Combine slides while preserving slide numbers.
        """

        sections = []

        for slide in slides:

            if not slide.get(
                "success",
                False,
            ):

                continue

            slide_number = slide.get(
                "slide_number"
            )

            title = slide.get(
                "title",
                "",
            )

            text = slide.get(
                "text",
                "",
            )

            if not text and not title:
                continue

            header = (
                f"[Slide {slide_number}]"
            )

            if title:

                header += (
                    f" - {title}"
                )

            sections.append(
                f"{header}\n{text}"
            )

        return "\n\n".join(
            sections
        )

    # =========================================================
    # RAG-READY SLIDE
    # =========================================================

    @staticmethod
    def create_rag_slide(
        slide: dict[str, Any],
    ) -> dict[str, Any]:
        """
        Convert processed slide into a RAG-friendly object.

        Embeddings and chunking happen later.
        """

        return {
            "slide_number": slide.get(
                "slide_number"
            ),
            "title": slide.get(
                "title",
                "",
            ),
            "text": slide.get(
                "text",
                "",
            ),
            "metadata": {
                "source_type": "pptx",
                "slide_number": slide.get(
                    "slide_number"
                ),
                "slide_type": slide.get(
                    "slide_type"
                ),
                "ocr_used": slide.get(
                    "ocr_used",
                    False,
                ),
                "image_count": len(
                    slide.get(
                        "images",
                        [],
                    )
                ),
                "table_count": len(
                    slide.get(
                        "tables",
                        [],
                    )
                ),
                "has_notes": bool(
                    slide.get(
                        "notes",
                        "",
                    )
                ),
                "structure": slide.get(
                    "structure",
                    {},
                ),
            },
        }

    # =========================================================
    # QUALITY VALIDATION
    # =========================================================

    def validate_quality(
        self,
        slides: list[dict[str, Any]],
        full_text: str,
    ) -> dict[str, Any]:
        """
        Validate presentation extraction quality.
        """

        total_slides = len(
            slides
        )

        successful_slides = sum(
            1
            for slide in slides
            if slide.get(
                "success",
                False,
            )
        )

        slides_with_text = sum(
            1
            for slide in slides
            if slide.get(
                "text",
                "",
            ).strip()
        )

        slides_with_ocr = sum(
            1
            for slide in slides
            if slide.get(
                "ocr_used",
                False,
            )
        )

        if total_slides == 0:

            extraction_percentage = 0.0

        else:

            extraction_percentage = round(
                (
                    slides_with_text
                    / total_slides
                )
                * 100,
                2,
            )

        if not full_text.strip():

            status = "poor"

        elif slides_with_text < total_slides:

            status = "partial"

        else:

            status = "good"

        return {
            "status": status,
            "total_slides": total_slides,
            "successful_slides": (
                successful_slides
            ),
            "slides_with_text": (
                slides_with_text
            ),
            "slides_with_ocr": (
                slides_with_ocr
            ),
            "text_length": len(
                full_text
            ),
            "extraction_percentage": (
                extraction_percentage
            ),
        }

    # =========================================================
    # TEXT CLEANING
    # =========================================================

    @staticmethod
    def clean_text(
        text: Optional[str],
    ) -> str:
        """
        Clean extracted presentation text.
        """

        if not text:

            return ""

        text = text.replace(
            "\r\n",
            "\n",
        )

        text = text.replace(
            "\r",
            "\n",
        )

        text = text.replace(
            "\x00",
            "",
        )

        text = re.sub(
            r"[ \t]+",
            " ",
            text,
        )

        text = re.sub(
            r" *\n *",
            "\n",
            text,
        )

        text = re.sub(
            r"\n{3,}",
            "\n\n",
            text,
        )

        return text.strip()

    # =========================================================
    # DOCUMENT ID
    # =========================================================

    @staticmethod
    def generate_document_id(
        path: Path,
    ) -> str:
        """
        Generate a deterministic document ID.
        """

        value = str(
            path
        ).encode(
            "utf-8"
        )

        return hashlib.sha256(
            value
        ).hexdigest()[:32]

    # =========================================================
    # SIMPLE TEXT EXTRACTION
    # =========================================================

    def extract_text(
        self,
        pptx_path: Union[str, Path],
    ) -> str:
        """
        Convenience method for extracting only
        presentation text.
        """

        result = self.process(
            pptx_path,
            use_ocr=False,
            include_notes=False,
        )

        if not result.get(
            "success",
            False,
        ):

            raise RuntimeError(
                result.get(
                    "error",
                    "PPTX processing failed.",
                )
            )

        return result.get(
            "full_text",
            "",
        )

    # =========================================================
    # SLIDE COUNT
    # =========================================================

    def get_slide_count(
        self,
        pptx_path: Union[str, Path],
    ) -> int:
        """
        Return the number of slides.
        """

        if Presentation is None:

            raise RuntimeError(
                "python-pptx is not installed."
            )

        path = Path(
            pptx_path
        ).resolve()

        if not path.exists():

            raise FileNotFoundError(
                path
            )

        presentation = Presentation(
            str(path)
        )

        return len(
            presentation.slides
        )

    # =========================================================
    # FAILURE RESPONSE
    # =========================================================

    @staticmethod
    def _failure(
        error: str,
    ) -> dict[str, Any]:
        """
        Standard failure response.
        """

        return {
            "success": False,
            "error": error,
        }


# =============================================================
# DEFAULT INSTANCE
# =============================================================

pptx_processor = PPTXProcessor()
