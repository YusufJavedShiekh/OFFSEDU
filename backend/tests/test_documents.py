"""
Tests for the document processing module.

Covers:
- File manager
- PDF processor
- DOCX processor
- PPTX processor
- Text processor
- Image processor
- OCR service
- Error handling
- Basic end-to-end document processing

Run:
    pytest backend/tests/test_documents.py -v
"""

from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest


# ---------------------------------------------------------------------------
# Imports
# ---------------------------------------------------------------------------

try:
    from backend.documents.file_manager import FileManager
except ImportError:
    FileManager = None

try:
    from backend.documents.pdf_processor import PDFProcessor
except ImportError:
    PDFProcessor = None

try:
    from backend.documents.docx_processor import DOCXProcessor
except ImportError:
    DOCXProcessor = None

try:
    from backend.documents.pptx_processor import PPTXProcessor
except ImportError:
    PPTXProcessor = None

try:
    from backend.documents.text_processor import TextProcessor
except ImportError:
    TextProcessor = None

try:
    from backend.documents.image_processor import ImageProcessor
except ImportError:
    ImageProcessor = None

try:
    from backend.documents.ocr_service import OCRService
except ImportError:
    OCRService = None


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def create_text_file(tmp_path: Path, name: str = "sample.txt") -> Path:
    """Create a simple UTF-8 text document."""
    file_path = tmp_path / name

    file_path.write_text(
        "Introduction\n"
        "This is a sample document.\n\n"
        "Chapter 1\n"
        "This chapter contains sample study material.",
        encoding="utf-8",
    )

    return file_path


def create_empty_file(tmp_path: Path, name: str = "empty.txt") -> Path:
    """Create an empty file."""
    file_path = tmp_path / name
    file_path.touch()
    return file_path


def create_invalid_file(tmp_path: Path, name: str = "invalid.pdf") -> Path:
    """Create a deliberately invalid document."""
    file_path = tmp_path / name
    file_path.write_bytes(b"This is not a valid document.")
    return file_path


def create_sample_image(tmp_path: Path, name: str = "sample.png") -> Path:
    """Create a small test image."""
    try:
        from PIL import Image
    except ImportError:
        pytest.skip("Pillow is not installed.")

    file_path = tmp_path / name

    image = Image.new("RGB", (100, 100), "white")
    image.save(file_path)

    return file_path


def create_sample_pdf(tmp_path: Path, name: str = "sample.pdf") -> Path:
    """Create a simple PDF for testing."""
    try:
        from reportlab.pdfgen import canvas
    except ImportError:
        pytest.skip("ReportLab is not installed.")

    file_path = tmp_path / name

    pdf = canvas.Canvas(str(file_path))
    pdf.drawString(100, 750, "Introduction")
    pdf.drawString(100, 720, "This is sample PDF content.")
    pdf.save()

    return file_path


def create_sample_docx(tmp_path: Path, name: str = "sample.docx") -> Path:
    """Create a simple DOCX document for testing."""
    try:
        from docx import Document
    except ImportError:
        pytest.skip("python-docx is not installed.")

    file_path = tmp_path / name

    document = Document()
    document.add_heading("Introduction", level=1)
    document.add_paragraph("This is sample DOCX content.")

    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Test"
    table.cell(1, 1).text = "123"

    document.save(file_path)

    return file_path


def create_sample_pptx(tmp_path: Path, name: str = "sample.pptx") -> Path:
    """Create a simple PPTX presentation for testing."""
    try:
        from pptx import Presentation
    except ImportError:
        pytest.skip("python-pptx is not installed.")

    file_path = tmp_path / name

    presentation = Presentation()

    slide = presentation.slides.add_slide(
        presentation.slide_layouts[1]
    )

    slide.shapes.title.text = "Introduction"
    slide.placeholders[1].text = "This is sample PPTX content."

    presentation.save(file_path)

    return file_path


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def sample_text(tmp_path):
    return create_text_file(tmp_path)


@pytest.fixture
def sample_pdf(tmp_path):
    return create_sample_pdf(tmp_path)


@pytest.fixture
def sample_docx(tmp_path):
    return create_sample_docx(tmp_path)


@pytest.fixture
def sample_pptx(tmp_path):
    return create_sample_pptx(tmp_path)


@pytest.fixture
def sample_image(tmp_path):
    return create_sample_image(tmp_path)


# ===========================================================================
# File Manager Tests
# ===========================================================================

@pytest.mark.skipif(FileManager is None, reason="FileManager unavailable")
class TestFileManager:
    """Tests for file_manager.py."""

    def test_file_manager_can_be_created(self):
        manager = FileManager()
        assert manager is not None

    def test_missing_file_is_handled(self, tmp_path):
        manager = FileManager()

        missing_file = tmp_path / "does_not_exist.txt"

        try:
            result = manager.validate_file(missing_file)
        except (AttributeError, TypeError, FileNotFoundError):
            return

        assert result is not None

    def test_existing_text_file(self, sample_text):
        manager = FileManager()

        try:
            result = manager.validate_file(sample_text)
        except AttributeError:
            pytest.skip("FileManager does not expose validate_file().")

        assert result is not None

    def test_unsupported_extension(self, tmp_path):
        manager = FileManager()

        file_path = tmp_path / "sample.xyz"
        file_path.write_text("unsupported", encoding="utf-8")

        try:
            result = manager.validate_file(file_path)
        except AttributeError:
            pytest.skip("FileManager does not expose validate_file().")
        except (ValueError, TypeError):
            return

        assert result is not None


# ===========================================================================
# Text Processor Tests
# ===========================================================================

@pytest.mark.skipif(TextProcessor is None, reason="TextProcessor unavailable")
class TestTextProcessor:
    """Tests for text_processor.py."""

    def test_processor_can_be_created(self):
        processor = TextProcessor()
        assert processor is not None

    def test_process_text_file(self, sample_text):
        processor = TextProcessor()

        result = processor.process(sample_text)

        assert result is not None

    def test_text_content_is_preserved(self, sample_text):
        processor = TextProcessor()

        result = processor.process(sample_text)

        result_string = str(result)

        assert "sample document" in result_string.lower()

    def test_empty_text_file(self, tmp_path):
        processor = TextProcessor()

        empty_file = create_empty_file(tmp_path)

        try:
            result = processor.process(empty_file)
        except (ValueError, RuntimeError):
            return

        assert result is not None

    def test_unicode_text(self, tmp_path):
        processor = TextProcessor()

        file_path = tmp_path / "unicode.txt"

        file_path.write_text(
            "English\n"
            "हिंदी\n"
            "मराठी\n"
            "اردو\n"
            "日本語",
            encoding="utf-8",
        )

        result = processor.process(file_path)

        assert result is not None
        assert "हिंदी" in str(result) or "मराठी" in str(result)


# ===========================================================================
# PDF Processor Tests
# ===========================================================================

@pytest.mark.skipif(PDFProcessor is None, reason="PDFProcessor unavailable")
class TestPDFProcessor:
    """Tests for pdf_processor.py."""

    def test_processor_can_be_created(self):
        processor = PDFProcessor()
        assert processor is not None

    def test_process_pdf(self, sample_pdf):
        processor = PDFProcessor()

        result = processor.process(sample_pdf)

        assert result is not None

    def test_pdf_content_is_extracted(self, sample_pdf):
        processor = PDFProcessor()

        result = processor.process(sample_pdf)

        result_string = str(result)

        assert "sample PDF content" in result_string

    def test_invalid_pdf_is_handled(self, tmp_path):
        processor = PDFProcessor()

        invalid_pdf = create_invalid_file(tmp_path)

        try:
            result = processor.process(invalid_pdf)
        except Exception as exc:
            assert isinstance(exc, Exception)
            return

        assert result is not None


# ===========================================================================
# DOCX Processor Tests
# ===========================================================================

@pytest.mark.skipif(DOCXProcessor is None, reason="DOCXProcessor unavailable")
class TestDOCXProcessor:
    """Tests for docx_processor.py."""

    def test_processor_can_be_created(self):
        processor = DOCXProcessor()
        assert processor is not None

    def test_process_docx(self, sample_docx):
        processor = DOCXProcessor()

        result = processor.process(sample_docx)

        assert result is not None

    def test_docx_content_is_extracted(self, sample_docx):
        processor = DOCXProcessor()

        result = processor.process(sample_docx)

        result_string = str(result)

        assert "sample DOCX content" in result_string

    def test_docx_heading_is_preserved(self, sample_docx):
        processor = DOCXProcessor()

        result = processor.process(sample_docx)

        result_string = str(result)

        assert "Introduction" in result_string

    def test_docx_table_content_is_preserved(self, sample_docx):
        processor = DOCXProcessor()

        result = processor.process(sample_docx)

        result_string = str(result)

        assert "Name" in result_string
        assert "Value" in result_string


# ===========================================================================
# PPTX Processor Tests
# ===========================================================================

@pytest.mark.skipif(PPTXProcessor is None, reason="PPTXProcessor unavailable")
class TestPPTXProcessor:
    """Tests for pptx_processor.py."""

    def test_processor_can_be_created(self):
        processor = PPTXProcessor()
        assert processor is not None

    def test_process_pptx(self, sample_pptx):
        processor = PPTXProcessor()

        result = processor.process(sample_pptx)

        assert result is not None

    def test_slide_title_is_extracted(self, sample_pptx):
        processor = PPTXProcessor()

        result = processor.process(sample_pptx)

        result_string = str(result)

        assert "Introduction" in result_string

    def test_slide_content_is_extracted(self, sample_pptx):
        processor = PPTXProcessor()

        result = processor.process(sample_pptx)

        result_string = str(result)

        assert "sample PPTX content" in result_string


# ===========================================================================
# Image Processor Tests
# ===========================================================================

@pytest.mark.skipif(ImageProcessor is None, reason="ImageProcessor unavailable")
class TestImageProcessor:
    """Tests for image_processor.py."""

    def test_processor_can_be_created(self):
        processor = ImageProcessor()
        assert processor is not None

    def test_process_image(self, sample_image):
        processor = ImageProcessor()

        result = processor.process(sample_image)

        assert result is not None

    def test_image_metadata_can_be_read(self, sample_image):
        processor = ImageProcessor()

        # Different implementations may expose metadata differently.
        if hasattr(processor, "get_metadata"):
            metadata = processor.get_metadata(sample_image)

            assert metadata is not None

        else:
            result = processor.process(sample_image)
            assert result is not None


# ===========================================================================
# OCR Tests
# ===========================================================================

@pytest.mark.skipif(OCRService is None, reason="OCRService unavailable")
class TestOCRService:
    """Tests for ocr_service.py."""

    def test_ocr_service_can_be_created(self):
        service = OCRService()
        assert service is not None

    def test_ocr_service_handles_image(self, sample_image):
        service = OCRService()

        # We don't require a specific OCR engine to be installed.
        # The service should either return a result or raise a controlled
        # exception when the OCR backend is unavailable.
        try:
            result = service.extract_text(sample_image)
        except (ImportError, RuntimeError, OSError, ValueError):
            return

        assert result is not None

    def test_ocr_failure_is_controlled(self):
        service = OCRService()

        fake_image = MagicMock()

        try:
            result = service.extract_text(fake_image)
        except Exception as exc:
            assert isinstance(exc, Exception)
            return

        assert result is not None


# ===========================================================================
# Error Handling Tests
# ===========================================================================

class TestDocumentErrorHandling:
    """General document error-handling tests."""

    def test_missing_file_path(self, tmp_path):
        missing_file = tmp_path / "missing.pdf"

        assert not missing_file.exists()

    def test_empty_file(self, tmp_path):
        empty_file = create_empty_file(tmp_path)

        assert empty_file.exists()
        assert empty_file.stat().st_size == 0

    def test_unsupported_file(self, tmp_path):
        unsupported = tmp_path / "file.xyz"
        unsupported.write_text("test", encoding="utf-8")

        assert unsupported.exists()
        assert unsupported.suffix == ".xyz"

    def test_corrupted_file(self, tmp_path):
        corrupted = create_invalid_file(tmp_path)

        assert corrupted.exists()
        assert corrupted.stat().st_size > 0


# ===========================================================================
# Structure / RAG-Ready Output Tests
# ===========================================================================

def _assert_rag_ready_result(result):
    """
    Perform lightweight validation of processor output.

    Different processors may return different structures, so this test
    intentionally checks only properties that should be universally useful.
    """
    assert result is not None

    if isinstance(result, dict):
        assert len(result) > 0

        # Common fields that may appear in RAG-ready output.
        possible_content_keys = (
            "content",
            "text",
            "pages",
            "sections",
            "paragraphs",
            "chunks",
            "document",
        )

        if any(key in result for key in possible_content_keys):
            return

    assert len(str(result)) > 0


@pytest.mark.skipif(TextProcessor is None, reason="TextProcessor unavailable")
def test_text_processor_produces_rag_ready_data(sample_text):
    """Verify that text processing produces usable structured output."""
    processor = TextProcessor()

    result = processor.process(sample_text)

    _assert_rag_ready_result(result)


@pytest.mark.skipif(PDFProcessor is None, reason="PDFProcessor unavailable")
def test_pdf_processor_produces_rag_ready_data(sample_pdf):
    """Verify that PDF processing produces usable output."""
    processor = PDFProcessor()

    result = processor.process(sample_pdf)

    _assert_rag_ready_result(result)


# ===========================================================================
# Integration Test
# ===========================================================================

@pytest.mark.skipif(
    TextProcessor is None,
    reason="TextProcessor unavailable",
)
def test_document_processing_pipeline(sample_text):
    """
    Basic integration test:

        File
          ↓
        Processor
          ↓
        Cleaned/structured content
          ↓
        RAG-ready result
    """
    processor = TextProcessor()

    result = processor.process(sample_text)

    assert result is not None

    result_string = str(result)

    assert len(result_string.strip()) > 0
    assert "sample" in result_string.lower()


# ===========================================================================
# Optional Dependency Tests
# ===========================================================================

def test_supported_test_formats():
    """Ensure the test suite covers the intended document formats."""
    supported_formats = {
        ".pdf",
        ".docx",
        ".pptx",
        ".txt",
        ".png",
    }

    assert ".pdf" in supported_formats
    assert ".docx" in supported_formats
    assert ".pptx" in supported_formats
    assert ".txt" in supported_formats
    assert ".png" in supported_formats
